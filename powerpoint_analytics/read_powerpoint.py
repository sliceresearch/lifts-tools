#!/usr/bin/env python
# coding: utf-8

# # Interactive Powerpoint Analytics Tool (IPAT)
# 
# A product of the **LIFTS** project (**Learning in Future Teaching Spaces**), 2019-2020.
# 
# Authors:
# * Mark Utting
# * Jacqui Blake

# In[89]:


from pptx import Presentation
import collections
from typing import List


# In[63]:


# pres_name = "ICT112_Week09_HTML.pptx"
# pres_name = "presentation-template_teaching-lecture_wide_2018.pptx"
pres_name = "presentation-example1.pptx"


# In[65]:


prs = Presentation(pres_name)


# In[66]:


print(pres_name, "has", len(prs.slides), "slides")


# ## Print summary of all slides, plus text

# In[67]:


# print some meta-data plus all text
n = 0
for slide in prs.slides:
    n += 1
    print("========== slide {} ========== [{}]".format(n, slide.slide_layout.name))
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        print(shape.name)
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                print("    " + run.text)


# ## Collect Layout Statistics

# In[90]:


def get_word_counts(slides) -> List[int]:
    """Count the amount of text in each slide."""
    word_count = []
    for slide in slides:
        # print(f"========== slide {len(text_count)+1} ========== [{slide.slide_layout.name}]")
        words = 0
        # find all text
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # print(shape.name)
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # print("    " + run.text)
                    words += len(run.text.split())
        word_count.append(words)
    return word_count


# In[91]:


def calculate_text_stars(word_counts) -> int:
    """Calculates a one-to-five star ranking for presentations that are not too text-heavy."""
    if word_counts == []:
        return 3
    words_per_slide = sum(word_counts) / len(word_counts)
    stars = 5 - abs(words_per_slide - 35) / 8
    # print(stars)
    return max(0, min(5, int(stars + 0.5)))

assert calculate_text_stars([]) == 3
assert calculate_text_stars([30, 40]) == 5
assert calculate_text_stars([10, 10]) == 2
assert calculate_text_stars([10, 30, 60]) == 5  # TODO: penalise text-heavy slides more?


# In[92]:


INTERACTIVE = set(["MS Forms Quiz/Survey", "Video", "Demo"])
layouts = collections.defaultdict(int)
layouts_interactive = 0
for slide in prs.slides:
    layouts[slide.slide_layout.name] += 1
    if slide.slide_layout.name in INTERACTIVE:
        layouts_interactive += 1


# In[94]:


# print the slide layout analytics
print(f"Slide Layout counts:")
for name, count in layouts.items():
    print(f"    {count} {name} slides.")
i_score = min(layouts_interactive, 5)
topic_score = ([1,1,3,5,5,4,3,2,1]+[1]*100)[layouts["Section Header"]]
word_count = get_word_counts(prs.slides)
# print("word counts:", word_count)
words_per_slide = sum(word_count) / len(word_count)
# ideal words/slide is 30-40 (5 stars)
text_stars = calculate_text_stars(word_count)
print(f"Interaction score   : {i_score} stars.")
print(f"Sections score      : {topic_score} stars.")
print(f"Accessibility score : coming soon...")
print(f"Text score          : {text_stars} stars   ({words_per_slide:.1f} words/slide)")
# Warn about very text-heavy slides
for slide, words in enumerate(word_count):
    if words > 65:
        print(f"    WARNING: slide {slide} has {words} words!")


# In[ ]:




